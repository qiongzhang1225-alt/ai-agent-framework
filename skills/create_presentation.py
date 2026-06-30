
def create_presentation(
    title: str,
    subtitle: str = "",
    author: str = "",
    slides: list = None,
    output_path: str = "presentation.pptx"
) -> str:
    """
    生成 PowerPoint 演示文稿。

    参数:
        title: 主标题
        subtitle: 副标题（可选）
        author: 作者信息（可选）
        slides: 幻灯片内容列表，每项为 {"heading": str, "body": str}
                其中 body 用换行符 \\n 分隔不同段落，用    缩进表示子层级
        output_path: 输出文件路径，默认 "presentation.pptx"

    返回:
        生成的文件路径

    示例:
        create_presentation(
            title="稳态平板法测导热系数",
            author="肖寒 物理24-1班",
            slides=[
                {"heading": "实验目的", "body": "了解热传导\\n学习稳态平板法"},
                {"heading": "实验结果", "body": "λ = 0.34 W/(m·K)"},
            ]
        )
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- 配色方案 ----
    C_DARK = RGBColor(0x1B, 0x3A, 0x5C)
    C_MID = RGBColor(0x2C, 0x5F, 0x8A)
    C_ACCENT = RGBColor(0x3A, 0x7C, 0xBD)
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    C_BLACK = RGBColor(0x33, 0x33, 0x33)
    C_GRAY = RGBColor(0x88, 0x88, 0x88)

    def _set_bg(slide, color):
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = color

    def _add_side_bar(slide):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.6), Inches(7.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_MID
        shape.line.fill.background()

    def _add_textbox(slide, text, left, top, width, height,
                     font_size=18, color=C_BLACK, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        lines = text.split('\n')
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line.lstrip()
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.alignment = align
            # 缩进表示子层级
            indent_level = len(line) - len(line.lstrip())
            p.level = indent_level // 2 if indent_level > 0 else 0
            p.space_after = Pt(6)
        return txBox

    def _add_accent_line(slide, left, top, width):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_ACCENT
        shape.line.fill.background()

    # ======== 第1页：封面 ========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_DARK)
    _add_side_bar(slide)

    # 标题
    _add_textbox(slide, title, Inches(1.2), Inches(1.8), Inches(10), Inches(1.5),
                 font_size=40, color=C_WHITE, bold=True)

    # 分隔线
    _add_accent_line(slide, Inches(1.2), Inches(3.6), Inches(4))

    # 副标题
    if subtitle:
        _add_textbox(slide, subtitle, Inches(1.2), Inches(4.0), Inches(10), Inches(0.7),
                     font_size=24, color=C_ACCENT, bold=False)

    # 作者
    if author:
        _add_textbox(slide, author, Inches(1.2), Inches(4.8), Inches(10), Inches(0.7),
                     font_size=18, color=C_ACCENT, bold=False)

    # ======== 内容页 ========
    if slides:
        for s in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(slide, C_WHITE)
            _add_side_bar(slide)

            heading = s.get("heading", "")
            body = s.get("body", "")

            # 标题
            _add_textbox(slide, heading, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
                         font_size=32, color=C_DARK, bold=True)
            _add_accent_line(slide, Inches(1), Inches(1.2), Inches(3))

            # 正文
            if body:
                # 检测是否有需要突出显示的结果（如 λ =）
                if "=" in body and ("λ" in body or "结果" in heading):
                    # 分离普通内容和结果行
                    lines = body.split('\n')
                    normal_lines = []
                    result_lines = []
                    for line in lines:
                        if '=' in line and ('λ' in line or 'W/(m' in line):
                            result_lines.append(line)
                        else:
                            normal_lines.append(line)

                    # 普通正文
                    if normal_lines:
                        _add_textbox(slide, '\n'.join(normal_lines),
                                     Inches(1), Inches(1.6), Inches(11), Inches(4.5),
                                     font_size=18, color=C_BLACK)

                    # 结果突出框
                    if result_lines:
                        result_text = '\n'.join(result_lines)
                        box = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(2.5), Inches(5.8), Inches(8), Inches(1.2)
                        )
                        box.fill.solid()
                        box.fill.fore_color.rgb = C_DARK
                        box.line.fill.background()
                        tf = box.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.text = result_text
                        p.font.size = Pt(28)
                        p.font.color.rgb = C_WHITE
                        p.font.bold = True
                        p.alignment = PP_ALIGN.CENTER
                else:
                    _add_textbox(slide, body, Inches(1), Inches(1.6), Inches(11), Inches(5.5),
                                 font_size=18, color=C_BLACK)

    # 保存
    prs.save(output_path)
    return output_path
